# --- Imports ---
import os

import logfire
from docx import Document
from docx.oxml.ns import qn
from docx.table import Table
from docx.text.paragraph import Paragraph
from pptx import Presentation


def _iter_block_items(doc):
    """Yield Paragraphs and Tables in TRUE document order.

    doc.paragraphs and doc.tables are two separate lists, so a table read from
    the second one has lost its position — it would land at the end of the text
    under whatever the last heading happened to be. Walking the body's XML
    children keeps each table under the heading it actually belongs to.
    """
    for child in doc.element.body.iterchildren():
        if child.tag == qn("w:p"):
            yield Paragraph(child, doc)
        elif child.tag == qn("w:tbl"):
            yield Table(child, doc)


def _heading_level(style_name: str) -> int:
    """"Heading 3" -> 3, anything else -> 0.

    Capped at 6: Word goes up to "Heading 9", markdown ATX stops at 6 and the
    chunker's HEADING_RE only matches #{1,6}, so a deeper level would be parsed
    as body text instead of a heading.
    """
    if not style_name.startswith("Heading "):
        return 0
    tail = style_name.split()[-1]
    return min(int(tail), 6) if tail.isdigit() else 0


def _parse_docx(file_path: str) -> str:
    """Word document -> markdown-ish text preserving the outline."""
    # STEP 1: doc = Document(file_path)

    # STEP 2: walk doc.paragraphs in order. For each non-empty paragraph, look at
    #         p.style.name:
    #
    #   STEP 2a: style starts with "Heading " -> emit "#" * level + " " + text
    #            level = int(style.name.split()[-1])
    #            Measured on cronjobs.docx: Heading 1x1, Heading 2x18,
    #            Heading 3x42, Heading 4x11 — a full 72-heading outline that
    #            the old loader flattened to plain text.
    #
    #   STEP 2b: style == "HTML Preformatted" -> this is code. There are 1,308 of
    #            these in cronjobs.docx and they come in RUNS of consecutive
    #            paragraphs, one per line. Buffer consecutive ones and emit the
    #            whole run as a single fenced block, the way parse_html handles
    #            <pre>. Emitting them one-per-block would produce 1,308 separate
    #            one-line blocks.
    #
    #   STEP 2c: anything else -> emit the text as a plain block

    # STEP 3: don't forget to flush a trailing code buffer after the loop ends

    # STEP 4: doc.tables — same treatment as parse_html: per row, join cells with
    #         " | "; join rows with "\n" so a split lands on a row boundary.
    #         NOTE: python-docx keeps tables OUT of doc.paragraphs, so they are
    #         invisible to STEP 2 and land at the end. That loses their position
    #         in the document. Decide whether you care; if you do, iterate
    #         doc.element.body children instead to get true document order.

    # STEP 5: return "\n\n".join(blocks)
    #
    # DECISION on the STEP 4 note: position DOES matter. A table dumped at the end
    # inherits the breadcrumb of the last heading in the file, which is simply the
    # wrong section. So this walks _iter_block_items() for true document order and
    # handles tables inline — STEP 4 happens inside the STEP 2 loop, not after it.
    doc = Document(file_path)

    blocks: list[str] = []
    code_buffer: list[str] = []

    def flush_code() -> None:
        """Emit the buffered run of preformatted lines as ONE fenced block."""
        if code_buffer:
            blocks.append("```\n" + "\n".join(code_buffer) + "\n```")
            code_buffer.clear()

    for item in _iter_block_items(doc):
        # STEP 4: tables — per row, join cells with " | "; join rows with "\n" so a
        #         split lands on a row boundary.
        if isinstance(item, Table):
            flush_code()
            rows = []
            for row in item.rows:
                cells = [cell.text.strip() for cell in row.cells]
                if any(cells):
                    rows.append(" | ".join(cells))
            if rows:
                blocks.append("\n".join(rows))
            continue

        style = item.style.name if item.style is not None else ""

        #   STEP 2b: style == "HTML Preformatted" -> this is code. Buffer the run.
        #            Blank preformatted lines are DROPPED, not kept. A blank line
        #            inside the fenced block would make _pack_paragraphs see two
        #            paragraphs (it splits on \n\s*\n), so the opening fence and the
        #            closing fence would land in different chunks — measured at 30
        #            broken blocks in cronjobs.docx before this line existed.
        if style == "HTML Preformatted":
            if item.text.strip():
                code_buffer.append(item.text.rstrip())
            continue

        # any non-preformatted paragraph ends the run
        flush_code()

        text = item.text.strip()
        if not text:
            continue

        #   STEP 2a: style starts with "Heading " -> emit "#" * level + " " + text
        level = _heading_level(style)
        if level:
            blocks.append("#" * level + " " + text)

        #   STEP 2c: anything else -> emit the text as a plain block
        else:
            blocks.append(text)

    # STEP 3: don't forget to flush a trailing code buffer after the loop ends
    flush_code()

    # STEP 5: return "\n\n".join(blocks)
    return "\n\n".join(blocks)


def _parse_pptx(file_path: str) -> str:
    """Slide deck -> markdown-ish text, one section per slide."""
    # STEP 1: prs = Presentation(file_path)

    # STEP 2: for each slide (enumerate from 1):

    #   STEP 2a: title = slide.shapes.title.text if there IS a title placeholder
    #            (58 of 66 slides have one; 8 do not — handle None).
    #            Fall back to the first non-empty text shape, else "Slide N".

    #   STEP 2b: emit the title as "## " + title — a slide IS a section, which is
    #            what gives every chunk from this deck a breadcrumb.

    #   STEP 2c: collect body text from the other shapes with a text frame.
    #            SKIP any shape whose text equals the title — measured: the title
    #            string appears on a SECOND shape on many slides, so without this
    #            check every title is emitted twice.

    #   STEP 2d: a slide with a title and no body is a section divider
    #            (measured: ~16 slides under 50 chars). Emitting the heading alone
    #            is correct — chunk_text drops empty-bodied sections, and the
    #            heading still shapes the breadcrumb of the slides that follow.

    # STEP 1: prs = Presentation(file_path)
    prs = Presentation(file_path)

    blocks: list[str] = []

    # STEP 2: for each slide (enumerate from 1):
    for number, slide in enumerate(prs.slides, start=1):
        title_shape = slide.shapes.title

        # every shape on the slide that actually carries text, in shape order
        texts = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            text = shape.text_frame.text.strip()
            if text:
                texts.append(text)

        #   STEP 2a: title = slide.shapes.title.text if there IS a title placeholder
        #            (58 of 66 slides have one; 8 do not — handle None).
        #            Fall back to the first non-empty text shape, else "Slide N".
        raw_title = title_shape.text.strip() if title_shape is not None else ""
        if not raw_title:
            raw_title = texts[0] if texts else f"Slide {number}"

        #   STEP 2b: emit the title as "## " + title — a slide IS a section, which is
        #            what gives every chunk from this deck a breadcrumb.
        #            Whitespace is collapsed first: 4 titles in architecture.pptx
        #            contain a newline, and HEADING_RE matches a single line only, so
        #            the tail of such a title would be read as body text instead.
        blocks.append("## " + " ".join(raw_title.split()))

        #   STEP 2c: collect body text from the other shapes with a text frame.
        #            SKIP any shape whose text equals the title — measured: the title
        #            string appears on a SECOND shape on many slides, so without this
        #            check every title is emitted twice.
        #            (On this deck it is every one of the 58 titled slides.)
        #   STEP 2d: a slide with a title and no body is a section divider. Emitting
        #            the heading alone is correct — chunk_text drops empty-bodied
        #            sections, and the heading still shapes the breadcrumb of the
        #            slides that follow.
        blocks.extend(text for text in texts if text != raw_title)

    # STEP 3: return "\n\n".join(blocks)
    return "\n\n".join(blocks)


def parse_office(file_path: str):
    """
    Parse .docx / .pptx into MARKDOWN-ish text that preserves structure.

    The previous version did "\\n".join(str(el) for el in partition(...)), which
    threw away every element's category AND produced no blank lines — so a
    66-slide deck arrived as one paragraph.
    """
    with logfire.span("📄 Office Document Parsing", filename=file_path):
        try:
            # STEP 1: dispatch on the extension -> _parse_docx / _parse_pptx,
            #         raise or warn on anything else
            ext = os.path.splitext(file_path)[1].lower()

            if ext == ".docx":
                text = _parse_docx(file_path)
            elif ext == ".pptx":
                text = _parse_pptx(file_path)
            else:
                # Raise rather than warn: the caller picked this loader by
                # extension, so reaching here means the routing is wrong and
                # returning "" would hide that as an empty document.
                raise ValueError(
                    f"parse_office got {ext!r} — expected .docx or .pptx. "
                    "Legacy .doc/.ppt are a different format and need conversion first."
                )

            # STEP 2: warn if the result is empty, log the char count, return
            if not text.strip():
                logfire.warning(
                    "Office document produced no text",
                    filename=file_path,
                    ext=ext,
                )

            logfire.info(
                "Parsed office document",
                filename=file_path,
                ext=ext,
                chars=len(text),
            )
            return text
        except Exception as e:
            logfire.error(f"❌ Office Parse Failed: {e}")
            raise e